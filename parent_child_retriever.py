from langchain_community.vectorstores import FAISS
import os
from dotenv import load_dotenv
import pandas as pd
import glob 
from langchain.docstore.document import Document
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
# from langchain_openai import ChatOpenAI
# from langchain_openai import OpenAIEmbeddings
from datetime import datetime
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from langchain.retrievers import ParentDocumentRetriever
from langchain.storage import InMemoryStore
from pptx import Presentation 
from langchain_community.document_loaders import UnstructuredExcelLoader
import networkx as nx
import pickle


load_dotenv()

skiprows = None
# skiprows = 1

# finds all csv and pdf files from the given directory
def find_files(directory):
    csv_files = glob.glob(f"{directory}/**/*.csv", recursive=True)
    pdf_files = glob.glob(f"{directory}/**/*.pdf", recursive=True)
    ppt_files = glob.glob(f"{directory}/**/*.pptx", recursive=True)
    excel_files = glob.glob(f"{directory}/**/*.xlsx", recursive=True)
    text_files = glob.glob(f"{directory}/**/*.txt", recursive=True)
    return csv_files, pdf_files, ppt_files, excel_files, text_files

# loads all csv data
def load_csv_data(csv_files):
    documents = []
    for file in csv_files:
        df = pd.read_csv(file,skiprows=skiprows)   # Read CSV file, skipping the first row (For Mirav Vista) because actual data starts from second row
        document = df.to_json(orient="records", lines=True)
        doc=Document(
                    page_content=document,
                    metadata={"source": file,"type":"csv"} 
                )
        documents.append(doc)
    return documents

# loads all pdf data
def load_pdf_data(pdf_files):
    documents = []
    for pdf in pdf_files:
        reader = PdfReader(pdf)
        for page in reader.pages:
            doc = Document(
                        page_content=page.extract_text(), 
                        metadata={"source":pdf,"type":"pdf"}
                        )
            documents.append(doc)
    return documents

def load_ppt_data(ppt_files):
    documents = []
    for ppt in ppt_files:
        prs = Presentation(ppt)
        for slide in prs.slides: 
            text = ""
            text = "\n".join(
                getattr(shape, "text", "") for shape in slide.shapes
            ) 
            doc = Document(
                        page_content=text,
                        metadata={"source":ppt,"type":"ppt"}
                        )
            documents.append(doc)
    return documents

def load_excel_data(excel_files):
    documents = []
    for file in excel_files:
        loader = UnstructuredExcelLoader(file,mode="elements")
        document = loader.load()
        documents.extend(document)
    return documents

def load_text_data(text_files):
    documents = []
    for file in text_files:
        with open(file, "r") as f:
            text = f.read().replace("\n", " ")
            doc = Document(
                        page_content=text,
                        metadata={"source":file,"type":"txt"}
                        )         
            documents.append(doc)
    return documents

# creating embeddings and retreiver 
def create_embeddings(text_chunks):
    # Initialize OpenAI embeddings with API key from environment
    embeddings = GoogleGenerativeAIEmbeddings(model='models/embedding-001',google_api_key=os.getenv("GOOGLE_API_KEY"))
    # Create FAISS vector store from documents
    vector_store = FAISS.from_documents(documents=text_chunks,embedding=embeddings)
    # Saving faiss_index 
    vector_store.save_local("faiss_index")
    with open("text_chunks.pkl", "wb") as f:
        pickle.dump(text_chunks, f)
        
# Creates the vector store and retriever
def create_csv_pdf_vectorstore(csv_files, pdf_files, ppt_files, excel_files, text_files):
    csv_documents = load_csv_data(csv_files)     # Load CSV data
    pdf_documents = load_pdf_data(pdf_files)     # Load PDF data   
    ppt_documents = load_ppt_data(ppt_files)     # Load PPT data
    excel_documents = load_excel_data(excel_files) # Load Excel data
    text_documents =  load_text_data(text_files)    # Load text data
    all_documents = csv_documents + pdf_documents + ppt_documents + excel_documents + text_documents
    create_embeddings(all_documents)

# Stores the results of a query
def store_results(convo_uuid, query_id, path, method, result, docs):
    data = [{ "conversation_id": convo_uuid,
         "query_id": query_id,
         "test_time": datetime.now().isoformat(timespec="seconds"),
         "query": result['question'], 
         "result": result["output_text"],
         "method":method,
         "prompt":result["prompt"],
         "source_documents": docs}]
    
    new_row_df = pd.DataFrame.from_dict(data)

    if not os.path.exists(path):
        new_row_df.to_csv(path, index=False)
    else:
        df = pd.read_csv(path)
        pd.concat([df, new_row_df]).to_csv(path, index=False)

# Creates a conversational chain for Q&A
def get_conversational_chain():
    prompt_template = """
    Answer the question in a detailed and descriptive manner based on the provided context. 
    Make sure to cover all relevant aspects comprehensively. 
    If the information is not available in the provided context, respond with 'I don't know about this
    Context:\n {context}?\n
    Question:\n {question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro",temperature=0.7,api_key=os.getenv("GOOGLE_API_KEY"))     # Initialize model
    prompt = PromptTemplate(template=prompt_template,input_variables=["context","question"])    # Create prompt template
    chain = load_qa_chain(model,chain_type="stuff",prompt=prompt)
    return chain,prompt

# Handles user input and retrieves relevant documents
def user_input(user_question, retriever, convo_id, query_id, chain, prompt):
    retrieved_docs  = retriever.invoke(user_question)    # Retrieve documents based on user question
    # Run chain with retrieved documents
    response = run_chain(convo_id, query_id, user_question, retrieved_docs, chain, prompt) 
    return response,retrieved_docs

def run_chain(convo_id, query_id, user_question, docs, chain, prompt):
    response = chain.invoke(
        {"input_documents":docs, "question":user_question, "prompt":prompt})    # Run the chain with documents and question
    return response

def load_files_and_vectorstore(directory):
    csv_files, pdf_files, ppt_files, excel_files, text_files = find_files(directory)
    create_csv_pdf_vectorstore(csv_files, pdf_files, ppt_files, excel_files, text_files)  

def load_faiss_retriever():
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    
    with open("text_chunks.pkl", "rb") as f:
        text_chunks = pickle.load(f)
    
    # Initialize in-memory store
    store = InMemoryStore()
    # Define text splitters for child and parent documents
    child_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    parent_splitter = RecursiveCharacterTextSplitter(chunk_size=3000, chunk_overlap=200)    
    # Create parent document retriever
    parent_document_retriever = ParentDocumentRetriever(
        vectorstore=vector_store,
        docstore=store,
        child_splitter=child_splitter,
        parent_splitter=parent_splitter,
        search_kwargs={"k":100}  
    )
    # Add documents to the retriever
    parent_document_retriever.add_documents(text_chunks)
    return parent_document_retriever

